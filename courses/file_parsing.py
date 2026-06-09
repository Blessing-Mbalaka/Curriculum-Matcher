from pathlib import Path
from dataclasses import dataclass
from zipfile import BadZipFile

from django.core.exceptions import ValidationError


SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF",
    ".docx": "Word",
    ".txt": "text",
    ".pptx": "PowerPoint",
}


class IgnoredUploadedFile(Exception):
    pass


@dataclass
class ParsedUploadContent:
    text: str
    ignored_files: list[str]

    @property
    def ignored_count(self):
        return len(self.ignored_files)


def parse_uploaded_files(files):
    sections = []
    ignored_files = []
    for uploaded_file in files:
        try:
            parsed_text = parse_uploaded_file(uploaded_file)
        except IgnoredUploadedFile:
            ignored_files.append(uploaded_file.name)
            continue
        if parsed_text:
            sections.append(f"--- {uploaded_file.name} ---\n{parsed_text}")
    return ParsedUploadContent("\n\n".join(sections).strip(), ignored_files)


def parse_uploaded_file(uploaded_file):
    extension = Path(uploaded_file.name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValidationError(f"{uploaded_file.name}: unsupported file type. Supported files: {supported}.")

    try:
        if extension == ".txt":
            return parse_text_file(uploaded_file)
        if extension == ".pdf":
            return parse_pdf_file(uploaded_file)
        if extension == ".docx":
            return parse_docx_file(uploaded_file)
        if extension == ".pptx":
            return parse_pptx_file(uploaded_file)
    except IgnoredUploadedFile:
        raise
    except ValidationError:
        raise
    except Exception as exc:
        raise ValidationError(f"{uploaded_file.name}: could not parse file ({exc}).") from exc

    return ""


def parse_text_file(uploaded_file):
    data = uploaded_file.read()
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding).strip()
        except UnicodeDecodeError:
            continue
    raise ValidationError(f"{uploaded_file.name}: could not decode text file.")


def parse_pdf_file(uploaded_file):
    try:
        from pypdf import PdfReader
        from pypdf.errors import DependencyError, FileNotDecryptedError, PdfReadError
    except ImportError as exc:
        raise ValidationError("PDF parsing requires the pypdf package. Install requirements first.") from exc

    try:
        reader = PdfReader(uploaded_file)
        if getattr(reader, "is_encrypted", False):
            decrypt_result = reader.decrypt("")
            if not decrypt_result:
                raise IgnoredUploadedFile(uploaded_file.name)
        text = "\n\n".join((page.extract_text() or "").strip() for page in reader.pages).strip()
    except DependencyError as exc:
        raise ValidationError(
            f"{uploaded_file.name}: this PDF uses AES encryption. "
            "Install cryptography in the virtual environment, then upload again."
        ) from exc
    except FileNotDecryptedError as exc:
        raise IgnoredUploadedFile(uploaded_file.name) from exc
    except PdfReadError as exc:
        raise ValidationError(f"{uploaded_file.name}: could not read PDF structure ({exc}).") from exc

    if not text:
        raise ValidationError(
            f"{uploaded_file.name}: no extractable text was found. "
            "If this is a scanned PDF, run OCR or upload a text/Word version."
        )
    return text


def parse_docx_file(uploaded_file):
    try:
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError as exc:
        raise ValidationError("Word parsing requires the python-docx package. Install requirements first.") from exc

    try:
        document = Document(uploaded_file)
    except (BadZipFile, PackageNotFoundError) as exc:
        raise IgnoredUploadedFile(uploaded_file.name) from exc
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    table_cells = []
    for table in document.tables:
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if values:
                table_cells.append(" | ".join(values))
    return "\n".join([*paragraphs, *table_cells]).strip()


def parse_pptx_file(uploaded_file):
    try:
        from pptx import Presentation
        from pptx.exc import PackageNotFoundError
    except ImportError as exc:
        raise ValidationError("PowerPoint parsing requires the python-pptx package. Install requirements first.") from exc

    try:
        presentation = Presentation(uploaded_file)
    except (BadZipFile, PackageNotFoundError) as exc:
        raise IgnoredUploadedFile(uploaded_file.name) from exc
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
        if text_parts:
            slides.append(f"Slide {index}\n" + "\n".join(text_parts))
    return "\n\n".join(slides).strip()
